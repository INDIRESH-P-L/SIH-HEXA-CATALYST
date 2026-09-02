"""Text extraction from uploaded documents.

PDF extraction uses **pdfplumber**, which is MIT licensed. PyMuPDF is
deliberately not used: it is AGPL-3.0, and many government IT policies prohibit
AGPL dependencies. That choice is recorded in the licence table in the README.

Page numbers are preserved through extraction because they travel all the way
to the generated question, letting a trainer check any item against its source
page.
"""

from __future__ import annotations

import io
from dataclasses import dataclass

from app.core.errors import ValidationFailedError
from app.core.logging import get_logger

log = get_logger(__name__)

#: Below this, a document has not really been read. Almost always a scanned PDF
#: with no text layer. Saying so is better than generating questions about
#: nothing (§11.1).
MIN_USEFUL_CHARS = 200

SUPPORTED_EXTENSIONS = {"pdf", "docx", "pptx"}

MIME_BY_EXTENSION = {
    "pdf": {"application/pdf"},
    "docx": {
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    },
    "pptx": {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    },
}


@dataclass(frozen=True)
class Page:
    """One page or slide of extracted text."""

    page_no: int
    text: str


@dataclass(frozen=True)
class Extraction:
    pages: list[Page]
    page_count: int
    char_count: int

    @property
    def full_text(self) -> str:
        return "\n\n".join(p.text for p in self.pages if p.text.strip())


def _extract_pdf(data: bytes) -> list[Page]:
    import pdfplumber

    pages: list[Page] = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for index, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            pages.append(Page(page_no=index, text=text))
    return pages


def _extract_docx(data: bytes) -> list[Page]:
    """DOCX has no intrinsic pages, so paragraphs are grouped into blocks.

    Tables are included: in training handouts they often carry the definitions
    worth asking about.
    """
    from docx import Document

    document = Document(io.BytesIO(data))
    parts: list[str] = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    if not parts:
        return []

    # ~40 paragraphs per synthetic page keeps page_no meaningful for review.
    block = 40
    return [
        Page(page_no=i + 1, text="\n".join(parts[i * block : (i + 1) * block]))
        for i in range((len(parts) + block - 1) // block)
    ]


def _extract_pptx(data: bytes) -> list[Page]:
    """One slide is one page."""
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    pages: list[Page] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
        pages.append(Page(page_no=index, text="\n".join(parts)))
    return pages


def extract(data: bytes, file_type: str) -> Extraction:
    """Extract text, preserving page numbers.

    Raises :class:`ValidationFailedError` when the document yields too little
    text to work with, so the interface can say why rather than producing
    nonsense.
    """
    extension = file_type.lower().lstrip(".")
    if extension not in SUPPORTED_EXTENSIONS:
        raise ValidationFailedError(
            f"Unsupported file type '{extension}'. Upload a PDF, DOCX or PPTX."
        )

    try:
        if extension == "pdf":
            pages = _extract_pdf(data)
        elif extension == "docx":
            pages = _extract_docx(data)
        else:
            pages = _extract_pptx(data)
    except ValidationFailedError:
        raise
    except Exception as exc:
        raise ValidationFailedError(
            f"Could not read this {extension.upper()} file: {exc}"
        ) from exc

    char_count = sum(len(p.text) for p in pages)
    if char_count < MIN_USEFUL_CHARS:
        raise ValidationFailedError(
            f"Only {char_count} characters of text could be extracted from this "
            f"{extension.upper()}. It is most likely a scanned document with no "
            "text layer. Optical character recognition is not part of this "
            "prototype; upload a text-based document instead."
        )

    log.info(
        "extracted %d pages, %d characters from a %s", len(pages), char_count, extension
    )
    return Extraction(pages=pages, page_count=len(pages), char_count=char_count)


def validate_upload(filename: str, content_type: str | None, size: int, limit: int) -> str:
    """Check an upload and return its normalised extension.

    The client filename is used only to read the extension. The stored object is
    named from a server-generated UUID (§13.7).
    """
    if size > limit:
        raise ValidationFailedError(
            f"File is {size / 1_048_576:.1f} MB; the limit is "
            f"{limit / 1_048_576:.0f} MB."
        )
    if size == 0:
        raise ValidationFailedError("The uploaded file is empty.")

    extension = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if extension not in SUPPORTED_EXTENSIONS:
        raise ValidationFailedError(
            f"Unsupported file extension '.{extension}'. Upload a PDF, DOCX or PPTX."
        )

    expected = MIME_BY_EXTENSION[extension]
    if content_type and content_type not in expected:
        # Browsers are inconsistent about these, so mismatch is a warning rather
        # than a rejection; the parser is the real gate.
        log.warning(
            "content-type %s does not match extension .%s", content_type, extension
        )
    return extension
